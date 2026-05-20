# -*- coding: utf-8 -*-
"""
pptx_pil_stub.py — 在 Anaconda PIL DLL 壞掉時，仍能用 python-pptx 產 PPTX。

用法：把本檔最上面那段 "PIL Stub" 複製到你的 build script 最前面，
      在 `from pptx import ...` 之前先 inject，即可繞過壞掉的 Pillow native DLL。

原理：python-pptx 在 import 時就會 `from PIL import Image / ImageFont / ImageDraw`；
      插入圖片時還會呼叫 PIL 讀取尺寸。我們提供一個最小 stub：
      - 用 struct 自己解 PNG(IHDR) / JPEG(SOFn) header 拿寬高
      - 補齊 python-pptx 會碰到的屬性（.size/.format/.mode/.info）

下半部是一個「文字 + 右側配圖」的最小 deck builder 範本，從 _slides_data.json 讀資料。
"""
import sys
import struct

# ===================== PIL STUB（複製這段到你的 build script 最前面）=====================
class _StubImage:
    info = {}
    def __init__(self, src):
        if hasattr(src, "read"):
            data = src.read()
        elif isinstance(src, (bytes, bytearray)):
            data = bytes(src)
        else:
            with open(src, "rb") as f:
                data = f.read()
        self._data = data
        self._parse()
    def _parse(self):
        d = self._data
        if d[:8] == b"\x89PNG\r\n\x1a\n":                       # PNG: IHDR
            self.size = (struct.unpack(">I", d[16:20])[0], struct.unpack(">I", d[20:24])[0])
            self.format, self.mode = "PNG", "RGBA"; return
        if d[:2] == b"\xff\xd8":                                # JPEG: walk to SOFn
            i = 2
            while i < len(d):
                if d[i] != 0xFF: i += 1; continue
                m = d[i + 1]
                if m == 0xD9: break
                if (0xC0 <= m <= 0xC3) or (0xC5 <= m <= 0xC7) or (0xC9 <= m <= 0xCB) or (0xCD <= m <= 0xCF):
                    self.size = (struct.unpack(">H", d[i+7:i+9])[0], struct.unpack(">H", d[i+5:i+7])[0])
                    self.format, self.mode = "JPEG", "RGB"; return
                i += 2 + struct.unpack(">H", d[i+2:i+4])[0]
        self.size, self.format, self.mode = (1024, 768), "PNG", "RGB"
    def close(self): pass
    def __enter__(self): return self
    def __exit__(self, *a): pass

def _stub_open(fp, mode="r"): return _StubImage(fp)

class _StubImageModule:
    open = staticmethod(_stub_open); Image = _StubImage
class _StubFont:
    def getsize(self, t): return (len(t) * 7, 14)
    def getbbox(self, t): return (0, 0, len(t) * 7, 14)
    def getlength(self, t): return len(t) * 7
class _StubImageFontModule:
    truetype = staticmethod(lambda *a, **k: _StubFont())
    load_default = staticmethod(lambda *a, **k: _StubFont())
    FreeTypeFont = _StubFont
class _StubImageDrawModule:
    Draw = staticmethod(lambda *a, **k: None)
class _StubPIL:
    Image = _StubImageModule; ImageFont = _StubImageFontModule; ImageDraw = _StubImageDrawModule
    UnidentifiedImageError = type("UnidentifiedImageError", (Exception,), {})

sys.modules["PIL"] = _StubPIL
sys.modules["PIL.Image"] = _StubImageModule
sys.modules["PIL.ImageFont"] = _StubImageFontModule
sys.modules["PIL.ImageDraw"] = _StubImageDrawModule
sys.modules["Image"] = _StubImageModule
# =================== PIL STUB 結束 ===================

# ===================== 最小 deck builder 範本 =====================
# 期望 _slides_data.json 為 list[ {section,kicker,title,subtitle,bullets[],image} ]
if __name__ == "__main__":
    import json
    from pathlib import Path
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    ROOT = Path(__file__).resolve().parent
    data_path = ROOT / "_slides_data.json"
    if not data_path.exists():
        print("找不到 _slides_data.json — 這是範本，請改成你的資料路徑"); sys.exit(0)

    slides = json.loads(data_path.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
    INK = RGBColor(0xF1,0xF5,0xF9); INK2 = RGBColor(0xCB,0xD5,0xE1); INK3 = RGBColor(0x94,0xA3,0xB8)
    BG = RGBColor(0x0E,0x11,0x16); EMER = RGBColor(0x10,0xB9,0x81); VIO = RGBColor(0xA7,0x8B,0xFA)
    blank = prs.slide_layouts[6]

    def tb(slide,l,t,w,h,txt,sz,col,bold=False,mono=False,align=None):
        b = slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=b.text_frame
        tf.word_wrap=True; p=tf.paragraphs[0]
        if align: p.alignment=align
        r=p.add_run(); r.text=txt
        r.font.name = "Consolas" if mono else "Microsoft JhengHei"
        r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col
        return b

    try:
        out_enc = sys.stdout
        out_enc.reconfigure(encoding="utf-8")
    except Exception:
        pass

    for i, s in enumerate(slides, 1):
        sl = prs.slides.add_slide(blank)
        sl.background.fill.solid(); sl.background.fill.fore_color.rgb = BG
        if s.get("section"): tb(sl,6.2,0.18,3.6,0.25,f"- {s['section']} -",10,EMER,mono=True,align=PP_ALIGN.RIGHT)
        if s.get("kicker"):  tb(sl,0.4,0.18,5.5,0.25,s["kicker"],11,VIO,mono=True)
        tb(sl,0.4,0.48,9.2,0.7,s["title"],28,INK,bold=True)
        if s.get("subtitle"): tb(sl,0.4,1.18,9.2,0.4,s["subtitle"],16,INK2)
        has_img=False
        ip=s.get("image")
        if ip and Path(ip).exists():
            try: sl.shapes.add_picture(ip,Inches(5.4),Inches(1.7),width=Inches(4.3),height=Inches(2.9)); has_img=True
            except Exception as e: print(f"[WARN] image {i}: {e}")
        bw = 4.8 if has_img else 9.2
        bx = sl.shapes.add_textbox(Inches(0.4),Inches(1.7),Inches(bw),Inches(3.4)); btf=bx.text_frame; btf.word_wrap=True
        for j,b in enumerate(s.get("bullets",[])):
            p = btf.paragraphs[0] if j==0 else btf.add_paragraph(); p.space_after=Pt(6)
            r=p.add_run(); r.text=f"> {b}"; r.font.name="Microsoft JhengHei"; r.font.size=Pt(13); r.font.color.rgb=INK2
        tb(sl,8.9,5.32,0.9,0.25,f"{i} / {len(slides)}",10,INK3,mono=True,align=PP_ALIGN.RIGHT)

    out = ROOT / "deck.pptx"
    try: prs.save(out)
    except PermissionError:
        out = ROOT / "deck_new.pptx"; prs.save(out); print("[NOTE] 原檔被鎖，改存 deck_new.pptx")
    print(f"[OK] {out} ({len(prs.slides)} slides)")
